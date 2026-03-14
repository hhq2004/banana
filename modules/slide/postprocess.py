from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import xml.etree.ElementTree as ET
import re
import html
import base64
from io import BytesIO

def folder_to_pptx(folder_path: str, output_path: str):
    """将文件夹中的多个 DrawIO XML 文件合并为一个 PPTX（每个 XML 一个 slide）"""
    
    folder = Path(folder_path)
    xml_files = sorted(folder.glob('*.drawio')) + sorted(folder.glob('*.xml'))
    
    if not xml_files:
        print(f"文件夹 {folder_path} 中没有找到 .drawio 或 .xml 文件")
        return
    
    # 读取第一个文件确定 slide 尺寸
    with open(xml_files[0], 'r', encoding='utf-8') as f:
        first_content = f.read()
    
    root = ET.fromstring(first_content)
    graph_model = root.find('.//mxGraphModel')
    page_width = float(graph_model.get('pageWidth', 1920))
    page_height = float(graph_model.get('pageHeight', 1080))
    
    # 创建 PPTX
    prs = Presentation()
    prs.slide_width = Emu(page_width * 9525)
    prs.slide_height = Emu(page_height * 9525)
    
    # 遍历每个 XML 文件，创建一个 slide
    for xml_file in xml_files:
        print(f"处理: {xml_file.name}")
        with open(xml_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        add_elements_to_slide(slide, content)  # 添加该 XML 的所有元素
    
    prs.save(output_path)
    print(f"已保存: {output_path}，共 {len(xml_files)} 页")

def drawio_to_pptx(drawio_content: str, output_path: str):
    """将 DrawIO XML 转换为 PPTX（支持文字和图片）"""
    
    # 解析 XML
    root = ET.fromstring(drawio_content)
    
    # 获取画布尺寸
    graph_model = root.find('.//mxGraphModel')
    page_width = int(graph_model.get('pageWidth', 4800))
    page_height = int(graph_model.get('pageHeight', 3584))
    
    # 创建 PPTX（自定义尺寸）
    prs = Presentation()
    prs.slide_width = Emu(page_width * 9525)
    prs.slide_height = Emu(page_height * 9525)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 遍历所有 mxCell
    for cell in root.iter('mxCell'):
        style = cell.get('style', '')
        geometry = cell.find('mxGeometry')
        
        if geometry is None:
            continue
        
        # 解析位置和尺寸
        x = float(geometry.get('x', 0))
        y = float(geometry.get('y', 0))
        w = float(geometry.get('width', 100))
        h = float(geometry.get('height', 50))
        
        left = Emu(x * 9525)
        top = Emu(y * 9525)
        width = Emu(w * 9525)
        height = Emu(h * 9525)
        
        # 处理图片元素
        if 'shape=image' in style:
            add_image_to_slide(slide, style, left, top, width, height)
        # 处理文字元素
        elif 'text' in style:
            value = cell.get('value', '')
            if value:
                add_text_to_slide(slide, value, style, left, top, width, height)
        # 处理基本形状（矩形、圆形等）
        else:
            add_shape_to_slide(slide, style, left, top, width, height)
    
    prs.save(output_path)
    print(f"已保存: {output_path}")


def add_image_to_slide(slide, style: str, left, top, width, height):
    """添加图片到幻灯片"""
    # 提取 base64 图片数据
    match = re.search(r'image=data:image/(\w+)[,;]([^"]+)', style)
    if not match:
        return
    
    img_format = match.group(1)  # png, jpg 等
    base64_data = match.group(2)
    
    # 移除可能的 base64 前缀
    if ';base64,' in base64_data:
        base64_data = base64_data.split(';base64,')[1]
    
    try:
        # 解码 base64
        img_bytes = base64.b64decode(base64_data)
        img_stream = BytesIO(img_bytes)
        
        # 添加图片到幻灯片
        slide.shapes.add_picture(img_stream, left, top, width, height)
    except Exception as e:
        print(f"图片添加失败: {e}")


def add_text_to_slide(slide, value: str, style: str, left, top, width, height):
    """添加文字到幻灯片"""
    # 解析样式
    font_size = extract_style_value(style, 'fontSize', 12)
    font_color = extract_style_value(style, 'fontColor', '#000000')
    font_family = extract_style_value(style, 'fontFamily', 'Arial')
    is_bold = 'fontStyle=1' in style or 'fontStyle=3' in style
    is_italic = 'fontStyle=2' in style or 'fontStyle=3' in style
    
    # 处理 HTML 实体和 LaTeX
    text = html.unescape(value)
    text = clean_latex(text)
    
    # 添加文本框
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = False
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_family
    run.font.bold = is_bold
    run.font.italic = is_italic
    run.font.color.rgb = hex_to_rgb(font_color)


def add_shape_to_slide(slide, style: str, left, top, width, height):
    """添加基本形状到幻灯片（可选实现）"""
    from pptx.enum.shapes import MSO_SHAPE
    
    # 解析填充和边框颜色
    fill_color = extract_style_value(style, 'fillColor', '#ffffff')
    stroke_color = extract_style_value(style, 'strokeColor', '#000000')
    
    # 根据样式选择形状类型
    if 'ellipse' in style:
        shape_type = MSO_SHAPE.OVAL
    elif 'rounded=1' in style:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape_type = MSO_SHAPE.RECTANGLE
    
    try:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        
        # 设置填充颜色
        if fill_color and fill_color != 'none':
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        
        # 设置边框颜色
        if stroke_color and stroke_color != 'none':
            shape.line.color.rgb = hex_to_rgb(stroke_color)
    except Exception as e:
        print(f"形状添加失败: {e}")


def extract_style_value(style: str, key: str, default):
    """从样式字符串中提取值"""
    pattern = rf'{key}=([^;]+)'
    match = re.search(pattern, style)
    if match:
        val = match.group(1)
        if isinstance(default, int):
            try:
                return int(val)
            except:
                return default
        return val
    return default


def clean_latex(text: str) -> str:
    """简化处理 LaTeX 公式"""
    text = re.sub(r'\\\(|\\\)', '', text)
    text = re.sub(r'\\mathbf\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\mathtt\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\mathsf\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\bf\s*', '', text)
    text = re.sub(r'\\left|\\right', '', text)
    text = re.sub(r'\\;|\\,|\\~', ' ', text)
    text = re.sub(r'\\frac\{([^}]*)\}\{([^}]*)\}', r'(\1)/(\2)', text)
    text = re.sub(r'\\underset\{[^}]*\}\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def hex_to_rgb(hex_color: str) -> RGBColor:
    """十六进制颜色转 RGB"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return RGBColor(0, 0, 0)
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


if __name__ == '__main__':
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python postprocess.py <input.drawio> [output.pptx]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else input_file.replace('.drawio', '.pptx')
    
    with open(input_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    drawio_to_pptx(content, output_file)